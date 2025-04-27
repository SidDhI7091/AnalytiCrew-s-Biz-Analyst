from pptx import Presentation

def parse_ppt(file):
    content = []
    try:
        ppt = Presentation(file.file)
        for slide in ppt.slides:
            slide_data = {
                'slide_number': slide.slide_id,
                'text': []
            }
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_data['text'].append(shape.text.strip())
            content.append(slide_data)
    except Exception as e:
        content = f"Error parsing PPT: {str(e)}"
    
    return content
